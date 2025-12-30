from __future__ import annotations

import io
from typing import List

from .session_store import PageData


class DocumentExtractionError(RuntimeError):
    pass


def _ext_from_filename(filename: str) -> str:
    name = (filename or "").lower().strip()
    if "." not in name:
        return ""
    return name.rsplit(".", 1)[-1]


def extract_pages(filename: str, content: bytes) -> tuple[str, List[PageData]]:
    """
    Extract text per page/slide/chunk for:
      - PDF  -> per page
      - PPTX -> per slide
      - DOCX -> chunked into "virtual pages"

    Returns:
      (doc_type, pages)

    doc_type is one of: "pdf" | "pptx" | "docx"

    NOTE: This file only provides extraction. Storage is handled by SessionStore.
    """
    ext = _ext_from_filename(filename)

    if ext == "pdf":
        return "pdf", _extract_pdf_pages(content)
    if ext == "pptx":
        return "pptx", _extract_pptx_slides(content)
    if ext == "docx":
        return "docx", _extract_docx_chunks(content)
    if ext in ["jpg", "jpeg", "png"]:
        # Images are treated as single-page docs with no extracted text (for now)
        # The Vision Tutor uses the raw file/image bytes, not the text.
        return "image", [PageData(index=0, text="")]

    raise DocumentExtractionError(f"Unsupported file type: .{ext or '?'}")


def _extract_pdf_pages(content: bytes) -> List[PageData]:
    try:
        from pypdf import PdfReader
    except Exception as e:
        raise DocumentExtractionError(
            "Missing dependency for PDF extraction. Install: pypdf"
        ) from e

    try:
        reader = PdfReader(io.BytesIO(content))
        pages: List[PageData] = []
        for idx, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            # normalize whitespace lightly
            text = " ".join(text.split())
            pages.append(PageData(index=idx, text=text))
        return pages
    except Exception as e:
        raise DocumentExtractionError(f"Failed to extract PDF text: {e}") from e


def _extract_pptx_slides(content: bytes) -> List[PageData]:
    try:
        from pptx import Presentation
    except Exception as e:
        raise DocumentExtractionError(
            "Missing dependency for PPTX extraction. Install: python-pptx"
        ) from e

    try:
        prs = Presentation(io.BytesIO(content))
        slides: List[PageData] = []
        for idx, slide in enumerate(prs.slides):
            texts: List[str] = []
            for shape in slide.shapes:
                # common text-containing shapes
                if hasattr(shape, "text") and shape.text:
                    texts.append(shape.text)
            joined = "\n".join(texts).strip()
            joined = " ".join(joined.split())
            slides.append(PageData(index=idx, text=joined))
        return slides
    except Exception as e:
        raise DocumentExtractionError(f"Failed to extract PPTX text: {e}") from e


def _extract_docx_chunks(content: bytes, chunk_chars: int = 2500) -> List[PageData]:
    """
    DOCX has no real pages. We chunk paragraphs into "virtual pages"
    using a simple char budget.
    """
    try:
        import docx  # python-docx
    except Exception as e:
        raise DocumentExtractionError(
            "Missing dependency for DOCX extraction. Install: python-docx"
        ) from e

    try:
        doc = docx.Document(io.BytesIO(content))
        paras = []
        for p in doc.paragraphs:
            t = (p.text or "").strip()
            if t:
                paras.append(t)

        chunks: List[str] = []
        buf: List[str] = []
        size = 0

        for t in paras:
            # +1 for newline/space join
            if size + len(t) + 1 > chunk_chars and buf:
                chunks.append("\n".join(buf))
                buf = []
                size = 0
            buf.append(t)
            size += len(t) + 1

        if buf:
            chunks.append("\n".join(buf))

        pages: List[PageData] = []
        for idx, ch in enumerate(chunks):
            ch = " ".join(ch.split())
            pages.append(PageData(index=idx, text=ch))

        return pages
    except Exception as e:
        raise DocumentExtractionError(f"Failed to extract DOCX text: {e}") from e